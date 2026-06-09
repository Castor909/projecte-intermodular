#!/usr/bin/env python3
"""
Generate VinylEth_Presentation.pptx — UD3 final presentation.
Run: python3 docs/build_presentation.py
Output: docs/VinylEth_Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (dark premium theme) ──────────────────────────────────────
BG     = RGBColor(0x18, 0x0E, 0x04)   # dark espresso background
CARD   = RGBColor(0x2A, 0x17, 0x0A)   # card / panel bg
CREAM  = RGBColor(0xF2, 0xEC, 0xE0)   # primary text
DIM    = RGBColor(0xA8, 0x9A, 0x88)   # secondary / muted text
ORANGE = RGBColor(0xD4, 0x62, 0x2A)   # burnt orange accent
GOLD   = RGBColor(0xE8, 0x8A, 0x50)   # lighter accent

DOCS = os.path.dirname(os.path.abspath(__file__))
SS   = os.path.join(DOCS, "screenshots")

# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def txt(slide, text, x, y, w, h,
        size=22, color=CREAM, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size      = Pt(size)
    r.font.color.rgb = color
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.name      = "Calibri"
    return box

def bullets(slide, items, x, y, w, h,
            size=18, color=CREAM, prefix="▸  "):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        r = p.add_run()
        r.text           = prefix + item
        r.font.size      = Pt(size)
        r.font.color.rgb = color
        r.font.name      = "Calibri"
    return box

def rect(slide, x, y, w, h, fill=CARD, line=False):
    shp = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp

def hrule(slide, x, y, w, color=ORANGE, h=0.025):
    r = rect(slide, x, y, w, h, fill=color)
    return r

def accent_bar(slide):
    """Left orange vertical bar — consistent across all slides."""
    rect(slide, 0, 0, 0.08, 7.5, fill=ORANGE)

def img(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        return slide.shapes.add_picture(
            path, Inches(x), Inches(y), Inches(w), Inches(h))
    return slide.shapes.add_picture(
        path, Inches(x), Inches(y), Inches(w))

def slide_header(slide, title):
    accent_bar(slide)
    txt(slide, title, 0.35, 0.25, 12.5, 0.72,
        size=34, color=ORANGE, bold=True)
    hrule(slide, 0.35, 1.0, 12.7)

# ── Slide builders ────────────────────────────────────────────────────────────

def s_title(prs):
    s = blank(prs); set_bg(s)
    accent_bar(s)

    txt(s, "VinylEth", 0.35, 1.2, 9, 1.7,
        size=82, color=ORANGE, bold=True)
    hrule(s, 0.35, 3.1, 7.5, color=CREAM, h=0.018)
    txt(s, "Informe Final del Proyecto", 0.35, 3.25, 9, 0.75,
        size=28, color=CREAM)
    txt(s, "Projecte Intermodular UD3  ·  CIFP Francesc de Borja Moll",
        0.35, 4.05, 10, 0.5, size=15, color=DIM)
    txt(s, "Stepan Andreev", 0.35, 5.2, 6, 0.5,
        size=20, color=CREAM, bold=True)
    txt(s, "Junio 2026  ·  github.com/Castor909/projecte-intermodular",
        0.35, 5.75, 10, 0.4, size=13, color=DIM)

    # Right info block
    for i, (label, value) in enumerate([
        ("Stack",    "MERN"),
        ("Pagos",    "Ethereum / MetaMask"),
        ("Red",      "Sepolia testnet"),
        ("APIs",     "Discogs · iTunes · Resend"),
    ]):
        y = 1.5 + i * 1.3
        rect(s, 9.8, y, 3.3, 1.1, fill=CARD)
        hrule(s, 9.8, y, 0.2, color=ORANGE, h=1.1)
        txt(s, label, 10.15, y + 0.08, 3.0, 0.38,
            size=12, color=DIM)
        txt(s, value, 10.15, y + 0.5, 3.0, 0.5,
            size=16, color=CREAM, bold=True)


def s_idea(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "El Proyecto")

    txt(s,
        "VinylEth es una tienda online de vinilos donde los clientes "
        "pagan con Ethereum a traves de MetaMask, sin necesidad de "
        "procesadores de pago tradicionales.",
        0.35, 1.2, 6.3, 1.3, size=19, color=CREAM)

    bullets(s, [
        "El vinilo ha vuelto: 45 millones de discos vendidos en 2023",
        "Nicho objetivo: amantes del vinilo + usuarios Web3",
        "Privacidad: sin tarjetas de credito, sin bancos",
        "Flujo e-commerce completo: catalogo → carrito → pago crypto",
        "Transaccion real en la red de pruebas Sepolia",
    ], 0.35, 2.65, 6.3, 3.5, size=17)

    img(s, os.path.join(SS, "featured_album_and_special_offers.png"),
        6.9, 1.15, 6.2)


def s_stack(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Tecnologias")

    cards = [
        ("React 19  +  React Router 7", "SPA frontend (Vite)",                        ORANGE),
        ("Node.js  +  Express 5",        "API REST backend",                           GOLD),
        ("MongoDB  +  Mongoose 9",       "Base de datos NoSQL",                        ORANGE),
        ("JWT  +  bcrypt",               "Autenticacion stateless · hash contrasenas", GOLD),
        ("MetaMask  (window.ethereum)",  "Pagos Ethereum — sin librerias extra",       ORANGE),
        ("Discogs API  +  iTunes API",   "Metadatos de albumes · previews de audio",   GOLD),
        ("Resend",                       "Recibos de pedido por correo",               ORANGE),
        ("PWA Manifest",                 "Instalable en movil",                        GOLD),
    ]

    for i, (title, desc, accent) in enumerate(cards):
        col = 0.35 + (i % 2) * 6.45
        row = 1.15 + (i // 2) * 1.45
        rect(s, col, row, 6.2, 1.28, fill=CARD)
        hrule(s, col, row, 0.22, color=accent, h=1.28)
        txt(s, title, col + 0.42, row + 0.1,  5.6, 0.5,
            size=18, color=CREAM, bold=True)
        txt(s, desc,  col + 0.42, row + 0.65, 5.6, 0.45,
            size=14, color=DIM)


def s_architecture(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Arquitectura")
    img(s, os.path.join(DOCS, "architecture.png"), 0.35, 1.15, 12.6)


def s_catalog(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Catalogo y Navegacion")

    img(s, os.path.join(SS, "search_catalogue.png"),
        0.35, 1.15, 7.8)

    bullets(s, [
        "Busqueda en tiempo real — debounce 250ms",
        "Filtro por genero obtenido de la base de datos",
        "8 opciones de ordenacion (precio, año, stock...)",
        "Paginacion con boton Cargar mas",
        "Badges de stock — Sin stock / Ultimas copias",
        "Skeleton loading cards",
        "Vistos recientemente (localStorage)",
    ], 8.4, 1.25, 4.7, 5.8, size=16)


def s_detail(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Detalle de Album y Vista Previa")

    img(s, os.path.join(SS, "album_preview.png"),
        0.35, 1.15, 7.8)

    bullets(s, [
        "Metadatos completos (sello, pais, formato)",
        "Tracklist con duraciones individuales",
        "Reproductor HTML5 personalizado",
        "  · Play/pausa, barra de progreso, buffering",
        "  · Preview de 30 segundos via iTunes",
        "Albumes similares (mismo genero, 4 items)",
        "Anadir al carrito / toggle de lista de deseos",
    ], 8.4, 1.25, 4.7, 5.8, size=16)


def s_payment(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Proceso de Pago con MetaMask")

    img(s, os.path.join(SS, "payment_process.png"),
        0.35, 1.15, 7.8)

    txt(s, "Flujo en 3 pasos:", 8.4, 1.25, 4.7, 0.45,
        size=16, color=GOLD, bold=True)

    steps = [
        "Carrito  →  Envio  →  Pago",
        "Direccion de envio prellenada desde perfil",
        "Cambio de red a Sepolia (0xaa36a7)",
        "eth_sendTransaction via MetaMask",
        "ETH → Wei con BigInt (sin errores de redondeo)",
        "Hash TX + enlace a Etherscan al confirmar",
        "Pedido guardado en BD  +  recibo por email",
        "Carrito vaciado automaticamente",
    ]
    bullets(s, steps, 8.4, 1.75, 4.7, 5.3,
            size=15, prefix="→  ")


def s_admin(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Panel de Administracion")

    img(s, os.path.join(SS, "admin_panel.png"),
        0.35, 1.15, 7.8)

    bullets(s, [
        "Dashboard: albumes, pedidos, ingresos, top ventas",
        "CRUD de albumes — formulario completo",
        "Importacion rapida de Discogs (ID → todos los campos)",
        "Autorellenado de preview de audio desde iTunes",
        "Operaciones en lote: descuento %, destacado, eliminar",
        "Exportacion CSV del catalogo de albumes",
        "Busqueda de pedidos por hash TX + rango de fechas",
        "Autenticacion separada por contrasena (x-admin-token)",
    ], 8.4, 1.25, 4.7, 5.8, size=15)


def s_problems(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Problemas y Soluciones")

    problems = [
        (
            "MongoDB se bloqueaba con SIGSEGV en la maquina de desarrollo",
            "mongodb-bin 8.2.7 nativo crasheaba al arrancar en Arch Linux.",
            "Docker mongo:7 — BD estable sin cambios en el codigo.",
        ),
        (
            "Perdida de precision al convertir ETH → Wei con coma flotante",
            "amount * 1e18 producia errores de redondeo para valores como 0.015 ETH.",
            "Aritmetica BigInt: separar en parte entera y decimal, escalar cada una.",
        ),
        (
            "ESLint: CartContext exportaba provider y hook en el mismo archivo",
            "El plugin React fast-refresh rechazaba las exportaciones mixtas.",
            "Separado en useCart.js + cartContextValue.js — mejora tambien la SoC.",
        ),
    ]

    for i, (title, problem, solution) in enumerate(problems):
        y = 1.2 + i * 2.0
        rect(s, 0.35, y, 12.7, 1.78, fill=CARD)
        hrule(s, 0.35, y, 0.22, color=ORANGE, h=1.78)
        txt(s, title,    0.75, y + 0.1,  12.0, 0.45,
            size=17, color=ORANGE, bold=True)
        txt(s, "⚠  " + problem,  0.75, y + 0.58, 12.0, 0.45,
            size=14, color=DIM)
        txt(s, "✓  " + solution, 0.75, y + 1.05, 12.0, 0.55,
            size=14, color=CREAM)


def s_evolution(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Evolucion del Proyecto")

    phases = [
        ("UD1A", "Configuracion Inicial", [
            "React + catalogo simulado",
            "Esqueleto de Express",
            "Sistema de diseno definido",
        ]),
        ("UD1B", "Backend + Carrito", [
            "MongoDB + modelo Album",
            "API REST (lista y detalle)",
            "Carrito con localStorage",
        ]),
        ("UD2A", "Wallet + Busqueda", [
            "Conexion MetaMask en header",
            "Controles de cantidad en carrito",
            "Busqueda, filtro y orden en catalogo",
        ]),
        ("UD3", "App Completa", [
            "Flujo de pago completo + pedidos",
            "Auth JWT + panel de admin",
            "Email, lista de deseos, PWA",
        ]),
    ]

    # Timeline bar
    hrule(s, 0.35, 4.3, 12.7, color=ORANGE, h=0.05)

    col_w = 3.15
    for i, (phase, title, items) in enumerate(phases):
        x = 0.35 + i * col_w

        # Phase label
        txt(s, phase, x, 1.15, col_w - 0.15, 0.7,
            size=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, title, x, 1.9, col_w - 0.15, 0.5,
            size=15, color=GOLD, align=PP_ALIGN.CENTER)
        hrule(s, x + 0.15, 2.48, col_w - 0.4, color=CARD, h=0.02)

        # Bullet items above timeline
        bullets(s, items, x + 0.15, 2.6, col_w - 0.3, 1.55,
                size=13, color=DIM, prefix="· ")

        # Dot on timeline
        dot_x = x + col_w / 2 - 0.2
        rect(s, dot_x, 4.14, 0.38, 0.38, fill=ORANGE)

        # Phase repeated below dot
        txt(s, phase, x, 4.65, col_w - 0.15, 0.5,
            size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def s_conclusions(prs):
    s = blank(prs); set_bg(s)
    slide_header(s, "Logros Obtenidos")

    # Stats row
    stats = [
        ("111", "funcionalidades"),
        ("17",  "endpoints API"),
        ("3",   "APIs externas"),
        ("4",   "fases de desarrollo"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 0.35 + i * 3.15
        rect(s, x, 1.2, 3.0, 1.5, fill=CARD)
        txt(s, num,   x, 1.25, 3.0, 0.9,
            size=54, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, label, x, 2.1,  3.0, 0.45,
            size=14, color=DIM, align=PP_ALIGN.CENTER)

    # Two columns
    txt(s, "Entregado", 0.35, 2.95, 6.2, 0.45,
        size=18, color=GOLD, bold=True)
    bullets(s, [
        "Flujo e-commerce completo: catalogo → carrito → pago crypto",
        "Pagos Ethereum via MetaMask en la red Sepolia",
        "Panel de administracion con importacion desde Discogs",
        "Auth JWT, historial de pedidos, recibos por email via Resend",
    ], 0.35, 3.48, 6.2, 2.8, size=16)

    txt(s, "Mejoras futuras", 6.75, 2.95, 6.2, 0.45,
        size=18, color=GOLD, bold=True)
    bullets(s, [
        "Contrato Solidity de custodia (proteccion al comprador)",
        "Reserva de stock durante el checkout",
        "Subida de imagenes via Cloudinary",
        "Suite de tests automatizados — Vitest + Jest + Supertest",
    ], 6.75, 3.48, 6.2, 2.8, size=16, color=DIM)


def s_thankyou(prs):
    s = blank(prs); set_bg(s)
    accent_bar(s)

    txt(s, "Gracias.", 0.35, 1.6, 10, 1.7,
        size=78, color=CREAM, bold=True)
    hrule(s, 0.35, 3.5, 8.5, color=ORANGE, h=0.025)
    txt(s, "Preguntas?", 0.35, 3.7, 8, 0.9,
        size=38, color=ORANGE)
    txt(s, "github.com/Castor909/projecte-intermodular",
        0.35, 5.1, 10, 0.5, size=18, color=DIM)
    txt(s, "Stepan Andreev  ·  Projecte Intermodular UD3  ·  Junio 2026",
        0.35, 5.7, 10, 0.4, size=13, color=DIM)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    s_title(prs)        # 1
    s_idea(prs)         # 2
    s_stack(prs)        # 3
    s_architecture(prs) # 4
    s_catalog(prs)      # 5
    s_detail(prs)       # 6
    s_payment(prs)      # 7
    s_admin(prs)        # 8
    s_problems(prs)     # 9
    s_evolution(prs)    # 10
    s_conclusions(prs)  # 11
    s_thankyou(prs)     # 12

    out = os.path.join(DOCS, "VinylEth_Presentation.pptx")
    prs.save(out)
    print(f"✓  Saved: {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
